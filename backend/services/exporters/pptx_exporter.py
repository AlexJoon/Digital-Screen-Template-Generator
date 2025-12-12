import io
import qrcode
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from .base import BaseExporter, ExportFormat, SlideData, TemplateStyle, TemplateConfig


class PPTXExporter(BaseExporter):
    """Export slides to PowerPoint format"""

    # Slide dimensions (16:9 widescreen)
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    @property
    def format(self) -> ExportFormat:
        return ExportFormat.PPTX

    @property
    def content_type(self) -> str:
        return "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    @property
    def file_extension(self) -> str:
        return ".pptx"

    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """Convert hex color to RgbColor"""
        hex_color = hex_color.lstrip('#')
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        return RGBColor(r, g, b)

    def _add_gradient_background(self, slide, template):
        """Add gradient background to slide"""
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.SLIDE_WIDTH, self.SLIDE_HEIGHT
        )
        background.line.fill.background()

        fill = background.fill
        fill.gradient()
        fill.gradient_angle = 135
        fill.gradient_stops[0].color.rgb = self._hex_to_rgb(template.background_color)
        fill.gradient_stops[1].color.rgb = self._hex_to_rgb(template.background_gradient_end)

        # Send to back
        spTree = slide.shapes._spTree
        sp = background._element
        spTree.remove(sp)
        spTree.insert(2, sp)

    def _add_text_box(
        self,
        slide,
        text: str,
        left: float,
        top: float,
        width: float,
        height: float,
        font_size: int,
        font_color: str,
        bold: bool = False,
        alignment: PP_ALIGN = PP_ALIGN.LEFT
    ):
        """Add a text box to the slide"""
        textbox = slide.shapes.add_textbox(
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        tf = textbox.text_frame
        tf.word_wrap = True
        tf.auto_size = None

        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = self._hex_to_rgb(font_color)
        p.font.bold = bold
        p.font.name = "Neue Haas Grotesk Display Pro"
        p.alignment = alignment

        return textbox

    def _add_circular_image(self, slide, image_data: bytes, left: float, top: float, size: float):
        """Add image to slide (circular cropping done via placeholder)"""
        image_stream = io.BytesIO(image_data)

        # Add the image
        picture = slide.shapes.add_picture(
            image_stream,
            Inches(left), Inches(top),
            Inches(size), Inches(size)
        )

        return picture

    def _generate_qr_code(self, url: str) -> bytes:
        """Generate QR code image from URL"""
        qr = qrcode.QRCode(
            version=1,
            error_correction=qrcode.constants.ERROR_CORRECT_L,
            box_size=10,
            border=2,
        )
        qr.add_data(url)
        qr.make(fit=True)

        img = qr.make_image(fill_color="black", back_color="white")

        # Convert to bytes
        output = io.BytesIO()
        img.save(output, format='PNG')
        output.seek(0)
        return output.getvalue()

    def _add_qr_code(self, slide, url: str, left: float, top: float, size: float):
        """Add QR code to slide"""
        qr_bytes = self._generate_qr_code(url)
        qr_stream = io.BytesIO(qr_bytes)

        picture = slide.shapes.add_picture(
            qr_stream,
            Inches(left), Inches(top),
            Inches(size), Inches(size)
        )

        return picture

    def _add_full_image_background(self, slide, image_data: bytes):
        """Add full-bleed background image with dark overlay"""
        image_stream = io.BytesIO(image_data)
        picture = slide.shapes.add_picture(
            image_stream,
            Inches(0), Inches(0),
            self.SLIDE_WIDTH, self.SLIDE_HEIGHT
        )
        # Send to back
        spTree = slide.shapes._spTree
        sp = picture._element
        spTree.remove(sp)
        spTree.insert(2, sp)

        # Add dark overlay
        overlay = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.SLIDE_WIDTH, self.SLIDE_HEIGHT
        )
        overlay.line.fill.background()
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
        # Note: python-pptx doesn't support transparency directly; using gradient as workaround
        overlay_fill = overlay.fill
        overlay_fill.gradient()
        overlay_fill.gradient_angle = 180
        overlay_fill.gradient_stops[0].color.rgb = RGBColor(0, 0, 0)
        overlay_fill.gradient_stops[0].position = 0.0
        overlay_fill.gradient_stops[1].color.rgb = RGBColor(0, 0, 0)
        overlay_fill.gradient_stops[1].position = 1.0

    def _add_rectangular_image(self, slide, image_data: bytes, left: float, top: float, width: float, height: float):
        """Add rectangular image to slide"""
        image_stream = io.BytesIO(image_data)
        picture = slide.shapes.add_picture(
            image_stream,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        return picture

    def _layout_full_hero(self, slide, slide_data: SlideData, template_style: TemplateStyle, template: TemplateConfig):
        """Full hero layout - large background image with text overlay at bottom"""
        # Add full background image
        if slide_data.image_data:
            self._add_full_image_background(slide, slide_data.image_data)
        else:
            self._add_gradient_background(slide, template)

        left_margin = 0.8
        content_width = 10.0

        # Caption at bottom
        if slide_data.caption:
            self._add_text_box(
                slide, slide_data.caption.upper(),
                left=left_margin, top=4.5, width=content_width, height=0.4,
                font_size=12, font_color=template.text_color, bold=False
            )

        # Headline
        self._add_text_box(
            slide, slide_data.headline,
            left=left_margin, top=5.0, width=content_width, height=1.2,
            font_size=36, font_color=template.text_color, bold=True
        )

        # Description
        if slide_data.description:
            self._add_text_box(
                slide, slide_data.description,
                left=left_margin, top=6.2, width=8.0, height=0.8,
                font_size=14, font_color=template.text_color, bold=False
            )

        # Author name (bottom left)
        if slide_data.author_name:
            self._add_text_box(
                slide, slide_data.author_name,
                left=left_margin, top=7.0, width=4.0, height=0.4,
                font_size=14, font_color=template.accent_color, bold=True
            )

        # QR code (bottom right)
        if slide_data.publication_link:
            self._add_qr_code(slide, slide_data.publication_link, left=11.5, top=5.8, size=1.2)
            self._add_text_box(
                slide, "Scan for more",
                left=11.2, top=7.0, width=1.5, height=0.3,
                font_size=9, font_color=template.text_color, bold=False, alignment=PP_ALIGN.CENTER
            )

    def _layout_split_text_primary(self, slide, slide_data: SlideData, template_style: TemplateStyle, template: TemplateConfig):
        """Split layout with text on left (2/3) and circular image on right (1/3)"""
        self._add_gradient_background(slide, template)

        left_margin = 0.6
        content_width = 8.0
        right_section_left = 9.5
        image_size = 3.0

        # Caption
        if slide_data.caption:
            self._add_text_box(
                slide, slide_data.caption.upper(),
                left=left_margin, top=0.5, width=content_width, height=0.4,
                font_size=12, font_color=template.text_color, bold=False
            )

        # Headline
        headline_top = 1.2 if slide_data.caption else 0.8
        self._add_text_box(
            slide, slide_data.headline,
            left=left_margin, top=headline_top, width=content_width, height=1.5,
            font_size=36, font_color=template.text_color, bold=True
        )

        # Description
        self._add_text_box(
            slide, slide_data.description,
            left=left_margin, top=headline_top + 1.8, width=content_width, height=2.5,
            font_size=18, font_color=template.text_color, bold=False
        )

        # Event details
        if slide_data.event_date or slide_data.event_time or slide_data.event_location:
            event_parts = []
            if slide_data.event_date:
                event_parts.append(f"Date: {slide_data.event_date}")
            if slide_data.event_time:
                event_parts.append(f"Time: {slide_data.event_time}")
            if slide_data.event_location:
                event_parts.append(f"Location: {slide_data.event_location}")
            event_text = "   |   ".join(event_parts)
            self._add_text_box(
                slide, event_text,
                left=left_margin, top=headline_top + 4.0, width=content_width, height=0.5,
                font_size=14, font_color=template.accent_color, bold=True
            )

        # Author name
        if slide_data.author_name:
            self._add_text_box(
                slide, slide_data.author_name,
                left=left_margin, top=6.0, width=content_width, height=0.5,
                font_size=16, font_color=template.accent_color, bold=True
            )

        # CBS branding
        self._add_text_box(
            slide, "Columbia Business School",
            left=left_margin, top=6.8, width=4.0, height=0.4,
            font_size=11, font_color=template.text_color, bold=False
        )

        # Circular image on right
        if slide_data.image_data:
            self._add_circular_image(slide, slide_data.image_data, left=right_section_left, top=2.0, size=image_size)

        # QR code
        if slide_data.publication_link:
            self._add_qr_code(slide, slide_data.publication_link, left=right_section_left + 0.9, top=5.5, size=1.2)
            self._add_text_box(
                slide, "Scan for more",
                left=right_section_left, top=6.75, width=image_size, height=0.3,
                font_size=10, font_color=template.text_color, bold=False, alignment=PP_ALIGN.CENTER
            )

    def _layout_split_image_primary(self, slide, slide_data: SlideData, template_style: TemplateStyle, template: TemplateConfig):
        """Split layout with large image on left (1/2) and text on right (1/2)"""
        # Add gradient background for right side
        self._add_gradient_background(slide, template)

        # Large image on left half
        if slide_data.image_data:
            self._add_rectangular_image(slide, slide_data.image_data, left=0, top=0, width=6.5, height=7.5)

        # Text content on right
        right_margin = 7.0
        content_width = 5.8

        # Caption
        if slide_data.caption:
            self._add_text_box(
                slide, slide_data.caption.upper(),
                left=right_margin, top=0.8, width=content_width, height=0.4,
                font_size=12, font_color=template.text_color, bold=False
            )

        # Headline
        headline_top = 1.5 if slide_data.caption else 1.0
        self._add_text_box(
            slide, slide_data.headline,
            left=right_margin, top=headline_top, width=content_width, height=1.5,
            font_size=30, font_color=template.text_color, bold=True
        )

        # Description
        self._add_text_box(
            slide, slide_data.description,
            left=right_margin, top=headline_top + 1.8, width=content_width, height=2.0,
            font_size=14, font_color=template.text_color, bold=False
        )

        # Event details
        if slide_data.event_date or slide_data.event_time or slide_data.event_location:
            event_parts = []
            if slide_data.event_date:
                event_parts.append(f"Date: {slide_data.event_date}")
            if slide_data.event_time:
                event_parts.append(f"Time: {slide_data.event_time}")
            if slide_data.event_location:
                event_parts.append(f"Location: {slide_data.event_location}")
            event_text = "\n".join(event_parts)
            self._add_text_box(
                slide, event_text,
                left=right_margin, top=headline_top + 3.8, width=content_width, height=1.0,
                font_size=12, font_color=template.accent_color, bold=False
            )

        # Author name
        if slide_data.author_name:
            self._add_text_box(
                slide, slide_data.author_name,
                left=right_margin, top=6.0, width=content_width, height=0.4,
                font_size=14, font_color=template.accent_color, bold=True
            )

        # CBS branding and QR
        self._add_text_box(
            slide, "Columbia Business School",
            left=right_margin, top=6.8, width=3.0, height=0.4,
            font_size=10, font_color=template.text_color, bold=False
        )

        if slide_data.publication_link:
            self._add_qr_code(slide, slide_data.publication_link, left=11.5, top=5.8, size=1.0)

    def _layout_circular_speaker(self, slide, slide_data: SlideData, template_style: TemplateStyle, template: TemplateConfig):
        """Centered circular image with text below - ideal for speaker highlights"""
        self._add_gradient_background(slide, template)

        center_x = 6.667  # Center of slide

        # Caption at top
        if slide_data.caption:
            self._add_text_box(
                slide, slide_data.caption.upper(),
                left=2.0, top=0.5, width=9.333, height=0.4,
                font_size=12, font_color=template.text_color, bold=False, alignment=PP_ALIGN.CENTER
            )

        # Centered circular image
        image_size = 3.0
        if slide_data.image_data:
            self._add_circular_image(slide, slide_data.image_data, left=center_x - image_size/2, top=1.0, size=image_size)

        # Author name (prominent)
        if slide_data.author_name:
            self._add_text_box(
                slide, slide_data.author_name,
                left=2.0, top=4.2, width=9.333, height=0.5,
                font_size=24, font_color=template.accent_color, bold=True, alignment=PP_ALIGN.CENTER
            )

        # Headline
        self._add_text_box(
            slide, slide_data.headline,
            left=1.5, top=4.8, width=10.333, height=1.0,
            font_size=28, font_color=template.text_color, bold=True, alignment=PP_ALIGN.CENTER
        )

        # Description
        if slide_data.description:
            self._add_text_box(
                slide, slide_data.description,
                left=2.0, top=5.9, width=9.333, height=0.8,
                font_size=14, font_color=template.text_color, bold=False, alignment=PP_ALIGN.CENTER
            )

        # Event details (centered)
        if slide_data.event_date or slide_data.event_time or slide_data.event_location:
            event_parts = []
            if slide_data.event_date:
                event_parts.append(slide_data.event_date)
            if slide_data.event_time:
                event_parts.append(slide_data.event_time)
            if slide_data.event_location:
                event_parts.append(slide_data.event_location)
            event_text = "  |  ".join(event_parts)
            self._add_text_box(
                slide, event_text,
                left=2.0, top=6.7, width=9.333, height=0.4,
                font_size=12, font_color=template.accent_color, bold=False, alignment=PP_ALIGN.CENTER
            )

        # CBS branding
        self._add_text_box(
            slide, "Columbia Business School",
            left=0.5, top=7.0, width=3.0, height=0.3,
            font_size=10, font_color=template.text_color, bold=False
        )

        # QR code (bottom right)
        if slide_data.publication_link:
            self._add_qr_code(slide, slide_data.publication_link, left=11.5, top=6.0, size=1.0)

    def _layout_text_only(self, slide, slide_data: SlideData, template_style: TemplateStyle, template: TemplateConfig):
        """Bold text-only layout with centered content"""
        self._add_gradient_background(slide, template)

        center_x = 6.667

        # Caption at top
        if slide_data.caption:
            self._add_text_box(
                slide, slide_data.caption.upper(),
                left=1.5, top=1.5, width=10.333, height=0.5,
                font_size=16, font_color=template.text_color, bold=False, alignment=PP_ALIGN.CENTER
            )

        # Large centered headline
        self._add_text_box(
            slide, slide_data.headline,
            left=1.0, top=2.5, width=11.333, height=2.0,
            font_size=48, font_color=template.text_color, bold=True, alignment=PP_ALIGN.CENTER
        )

        # Description
        if slide_data.description:
            self._add_text_box(
                slide, slide_data.description,
                left=2.0, top=4.8, width=9.333, height=1.2,
                font_size=18, font_color=template.text_color, bold=False, alignment=PP_ALIGN.CENTER
            )

        # Author name
        if slide_data.author_name:
            self._add_text_box(
                slide, slide_data.author_name,
                left=2.0, top=6.2, width=9.333, height=0.5,
                font_size=16, font_color=template.accent_color, bold=True, alignment=PP_ALIGN.CENTER
            )

        # CBS branding
        self._add_text_box(
            slide, "Columbia Business School",
            left=4.5, top=6.9, width=4.333, height=0.3,
            font_size=11, font_color=template.text_color, bold=False, alignment=PP_ALIGN.CENTER
        )

        # QR code (bottom right)
        if slide_data.publication_link:
            self._add_qr_code(slide, slide_data.publication_link, left=11.5, top=6.0, size=1.0)

    def _layout_media_vertical(self, slide, slide_data: SlideData, template_style: TemplateStyle, template: TemplateConfig):
        """Media mention layout with article card style"""
        self._add_gradient_background(slide, template)

        # Article card on left (white background)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(0.5),
            Inches(7.5), Inches(6.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.fill.background()

        # Image at top of card
        if slide_data.image_data:
            self._add_rectangular_image(slide, slide_data.image_data, left=0.6, top=0.6, width=7.3, height=3.0)

        # Article content in card
        card_text_color = "#181a1c"

        if slide_data.caption:
            self._add_text_box(
                slide, slide_data.caption.upper(),
                left=0.8, top=3.8, width=7.0, height=0.3,
                font_size=10, font_color=card_text_color, bold=False
            )

        self._add_text_box(
            slide, slide_data.headline,
            left=0.8, top=4.2, width=7.0, height=1.2,
            font_size=22, font_color=card_text_color, bold=True
        )

        if slide_data.description:
            self._add_text_box(
                slide, slide_data.description,
                left=0.8, top=5.5, width=7.0, height=1.2,
                font_size=12, font_color=card_text_color, bold=False
            )

        # Right side - Featured info
        if slide_data.author_name:
            self._add_text_box(
                slide, f"Featured: {slide_data.author_name}",
                left=8.5, top=1.0, width=4.3, height=0.5,
                font_size=16, font_color=template.accent_color, bold=True
            )

        # CBS branding
        self._add_text_box(
            slide, "Columbia Business School",
            left=8.5, top=6.5, width=4.0, height=0.4,
            font_size=11, font_color=template.text_color, bold=False
        )

        # QR code
        if slide_data.publication_link:
            self._add_qr_code(slide, slide_data.publication_link, left=9.5, top=4.0, size=1.5)
            self._add_text_box(
                slide, "Read Article",
                left=9.0, top=5.6, width=2.5, height=0.3,
                font_size=10, font_color=template.text_color, bold=False, alignment=PP_ALIGN.CENTER
            )

    def _layout_media_wide(self, slide, slide_data: SlideData, template_style: TemplateStyle, template: TemplateConfig):
        """Wide media layout with image on left"""
        self._add_gradient_background(slide, template)

        # Large image on left
        if slide_data.image_data:
            self._add_rectangular_image(slide, slide_data.image_data, left=0, top=0, width=5.5, height=7.5)

        # Content on right
        right_margin = 6.0
        content_width = 6.8

        if slide_data.caption:
            self._add_text_box(
                slide, slide_data.caption.upper(),
                left=right_margin, top=1.0, width=content_width, height=0.4,
                font_size=12, font_color=template.accent_color, bold=False
            )

        self._add_text_box(
            slide, slide_data.headline,
            left=right_margin, top=1.5, width=content_width, height=1.5,
            font_size=28, font_color=template.text_color, bold=True
        )

        if slide_data.description:
            self._add_text_box(
                slide, slide_data.description,
                left=right_margin, top=3.2, width=content_width, height=2.0,
                font_size=14, font_color=template.text_color, bold=False
            )

        if slide_data.author_name:
            self._add_text_box(
                slide, slide_data.author_name,
                left=right_margin, top=5.5, width=content_width, height=0.4,
                font_size=14, font_color=template.accent_color, bold=True
            )

        self._add_text_box(
            slide, "Columbia Business School",
            left=right_margin, top=6.8, width=3.0, height=0.4,
            font_size=10, font_color=template.text_color, bold=False
        )

        if slide_data.publication_link:
            self._add_qr_code(slide, slide_data.publication_link, left=11.5, top=5.8, size=1.2)

    def _layout_congrats_framed(self, slide, slide_data: SlideData, template_style: TemplateStyle, template: TemplateConfig):
        """Congratulations layout with framed image"""
        self._add_gradient_background(slide, template)

        # Framed image on left
        if slide_data.image_data:
            # Add decorative frame
            frame = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.8), Inches(1.5),
                Inches(4.5), Inches(4.5)
            )
            frame.fill.background()
            frame.line.color.rgb = self._hex_to_rgb(template.accent_color)
            frame.line.width = Pt(4)

            # Image inside frame
            self._add_rectangular_image(slide, slide_data.image_data, left=1.0, top=1.7, width=4.1, height=4.1)

        # Content on right
        right_margin = 5.8
        content_width = 6.8

        # "Congratulations" caption
        self._add_text_box(
            slide, (slide_data.caption or "CONGRATULATIONS").upper(),
            left=right_margin, top=1.5, width=content_width, height=0.5,
            font_size=18, font_color=template.accent_color, bold=False
        )

        # Honoree name (prominent)
        if slide_data.author_name:
            self._add_text_box(
                slide, slide_data.author_name,
                left=right_margin, top=2.2, width=content_width, height=1.0,
                font_size=32, font_color=template.text_color, bold=True
            )

        # Achievement headline
        self._add_text_box(
            slide, slide_data.headline,
            left=right_margin, top=3.5, width=content_width, height=1.2,
            font_size=22, font_color=template.accent_color, bold=True
        )

        # Description
        if slide_data.description:
            self._add_text_box(
                slide, slide_data.description,
                left=right_margin, top=4.8, width=content_width, height=1.5,
                font_size=14, font_color=template.text_color, bold=False
            )

        # CBS branding
        self._add_text_box(
            slide, "Columbia Business School",
            left=right_margin, top=6.8, width=4.0, height=0.4,
            font_size=11, font_color=template.text_color, bold=False
        )

    def _layout_podcast(self, slide, slide_data: SlideData, template_style: TemplateStyle, template: TemplateConfig):
        """Podcast layout with artwork on left"""
        self._add_gradient_background(slide, template)

        # Podcast artwork on left (square)
        if slide_data.image_data:
            self._add_rectangular_image(slide, slide_data.image_data, left=0.8, top=1.5, width=4.5, height=4.5)
        else:
            # Placeholder for podcast artwork
            placeholder = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.8), Inches(1.5),
                Inches(4.5), Inches(4.5)
            )
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = RGBColor(50, 50, 50)
            placeholder.line.fill.background()

        # Content on right
        right_margin = 5.8
        content_width = 6.8

        # Episode caption
        self._add_text_box(
            slide, (slide_data.caption or "PODCAST EPISODE").upper(),
            left=right_margin, top=1.5, width=content_width, height=0.4,
            font_size=12, font_color=template.text_color, bold=False
        )

        # Episode title
        self._add_text_box(
            slide, slide_data.headline,
            left=right_margin, top=2.0, width=content_width, height=1.5,
            font_size=26, font_color=template.text_color, bold=True
        )

        # Description
        if slide_data.description:
            self._add_text_box(
                slide, slide_data.description,
                left=right_margin, top=3.7, width=content_width, height=1.8,
                font_size=14, font_color=template.text_color, bold=False
            )

        # Host name
        if slide_data.author_name:
            self._add_text_box(
                slide, f"Host: {slide_data.author_name}",
                left=right_margin, top=5.8, width=content_width, height=0.4,
                font_size=14, font_color=template.accent_color, bold=True
            )

        # CBS branding
        self._add_text_box(
            slide, "Columbia Business School",
            left=right_margin, top=6.8, width=3.0, height=0.4,
            font_size=10, font_color=template.text_color, bold=False
        )

        # QR code for podcast link
        if slide_data.publication_link:
            self._add_qr_code(slide, slide_data.publication_link, left=11.5, top=5.5, size=1.2)
            self._add_text_box(
                slide, "Listen Now",
                left=11.2, top=6.75, width=1.8, height=0.3,
                font_size=10, font_color=template.text_color, bold=False, alignment=PP_ALIGN.CENTER
            )

    async def export(self, slide_data: SlideData) -> bytes:
        """Generate PowerPoint presentation from slide data"""
        template = slide_data.template
        template_style = slide_data.template_style

        # Create presentation with 16:9 aspect ratio
        prs = Presentation()
        prs.slide_width = self.SLIDE_WIDTH
        prs.slide_height = self.SLIDE_HEIGHT

        # Add blank slide
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)

        # Select layout based on template style
        if template_style:
            layout_type = template_style.layout_type
            image_position = template_style.image_position

            if layout_type == "full_hero" or image_position == "full":
                self._layout_full_hero(slide, slide_data, template_style, template)
            elif layout_type == "split_image_primary" or (image_position == "left" and layout_type not in ["media_wide", "podcast_standard", "podcast_feature"]):
                self._layout_split_image_primary(slide, slide_data, template_style, template)
            elif layout_type == "event_speaker" or image_position in ["circular", "center"]:
                self._layout_circular_speaker(slide, slide_data, template_style, template)
            elif image_position == "none":
                self._layout_text_only(slide, slide_data, template_style, template)
            elif layout_type == "media_vertical":
                self._layout_media_vertical(slide, slide_data, template_style, template)
            elif layout_type == "media_wide":
                self._layout_media_wide(slide, slide_data, template_style, template)
            elif layout_type == "congrats_framed":
                self._layout_congrats_framed(slide, slide_data, template_style, template)
            elif layout_type in ["podcast_standard", "podcast_feature"]:
                self._layout_podcast(slide, slide_data, template_style, template)
            else:
                # Default to split text primary
                self._layout_split_text_primary(slide, slide_data, template_style, template)
        else:
            # Fallback: use default split text primary layout
            self._layout_split_text_primary(slide, slide_data, template_style, template)

        # Save to bytes
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)

        return output.getvalue()
