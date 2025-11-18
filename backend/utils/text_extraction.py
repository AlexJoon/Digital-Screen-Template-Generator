"""
Text extraction utilities for various file formats.
"""
import io
from pypdf import PdfReader
from docx import Document
from pptx import Presentation


def extract_text_from_pdf(file_content: bytes) -> str:
    """Extract text from PDF file."""
    try:
        pdf_file = io.BytesIO(file_content)
        pdf_reader = PdfReader(pdf_file)

        text_parts = []
        for page in pdf_reader.pages:
            text_parts.append(page.extract_text())

        return "\n".join(text_parts)
    except Exception as e:
        raise Exception(f"Error extracting text from PDF: {str(e)}")


def extract_text_from_docx(file_content: bytes) -> str:
    """Extract text from Word document."""
    try:
        docx_file = io.BytesIO(file_content)
        doc = Document(docx_file)

        text_parts = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)

        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text_parts.append(cell.text)

        return "\n".join(text_parts)
    except Exception as e:
        raise Exception(f"Error extracting text from Word document: {str(e)}")


def extract_text_from_pptx(file_content: bytes) -> str:
    """Extract text from PowerPoint presentation."""
    try:
        pptx_file = io.BytesIO(file_content)
        prs = Presentation(pptx_file)

        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text.strip():
                        text_parts.append(shape.text)

        return "\n".join(text_parts)
    except Exception as e:
        raise Exception(f"Error extracting text from PowerPoint: {str(e)}")


def extract_text_from_file(file_content: bytes, content_type: str, filename: str) -> str:
    """
    Extract text from uploaded file based on content type.

    Args:
        file_content: Raw file bytes
        content_type: MIME type of the file
        filename: Original filename

    Returns:
        Extracted text content
    """
    if content_type == "text/plain":
        return file_content.decode("utf-8")

    elif content_type == "application/pdf":
        return extract_text_from_pdf(file_content)

    elif content_type in [
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword"
    ]:
        return extract_text_from_docx(file_content)

    elif content_type in [
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint"
    ]:
        return extract_text_from_pptx(file_content)

    else:
        raise ValueError(f"Unsupported file type: {content_type}")
